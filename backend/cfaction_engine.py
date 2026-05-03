"""
cfaction_engine — fonctions PURES pour :
- Analyser des pièces jointes (PDF/DOCX/XLSX/PPTX/SQLite/Image)
- Construire des fichiers téléchargeables (DOCX/PDF/XLSX/PPTX/plain)
- Exécuter du code Python en sandbox sécurisé

Ce module ne dépend PAS de MongoDB ni de FastAPI. Les fonctions retournent
des `bytes` (ou des dicts plats) et peuvent être testées/réutilisées partout.
La persistance en DB reste dans server.py (_persist_generated / _store_generated).
"""
from __future__ import annotations

import asyncio
import base64
import io
import logging
import os
import shutil
import sys
import tempfile
import time
from typing import Any, Dict, List, Optional

logger = logging.getLogger(__name__)


# ----------------------------------------------------------------- filename
def sanitize_filename(name: str, ext: str = "") -> str:
    """Nettoie un nom de fichier en conservant l'extension si présente."""
    stem = name or "file"
    extpart = ""
    if "." in stem:
        stem, _dot, extpart = stem.rpartition(".")
        extpart = f".{extpart}"
    base = "".join(c if (c.isalnum() or c in ("-", "_", " ")) else "_" for c in stem).strip()
    base = base.replace(" ", "_")[:80] or "file"
    return f"{base}{extpart}{ext}"


# ----------------------------------------------------------------- analyzers
def analyze_pdf(data: bytes) -> str:
    """Extrait le texte d'un PDF (max 40 pages, 30k caractères)."""
    try:
        from pypdf import PdfReader
        reader = PdfReader(io.BytesIO(data))
        parts = []
        for i, page in enumerate(reader.pages[:40]):
            try:
                t = page.extract_text() or ""
            except Exception:
                t = ""
            if t.strip():
                parts.append(f"[Page {i + 1}]\n{t.strip()}")
        return "\n\n".join(parts)[:30000]
    except Exception as e:
        logger.warning(f"PDF analyze failed: {e}")
        return ""


def analyze_docx(data: bytes) -> str:
    """Extrait le texte d'un DOCX."""
    try:
        from docx import Document
        doc = Document(io.BytesIO(data))
        parts = [p.text for p in doc.paragraphs if p.text and p.text.strip()]
        return "\n".join(parts)[:30000]
    except Exception as e:
        logger.warning(f"DOCX analyze failed: {e}")
        return ""


def analyze_xlsx(data: bytes) -> str:
    """Extrait le contenu XLSX (6 feuilles × 200 lignes max, format | séparé)."""
    try:
        from openpyxl import load_workbook
        wb = load_workbook(io.BytesIO(data), data_only=True, read_only=True)
        parts = []
        for ws in wb.worksheets[:6]:
            parts.append(f"=== Feuille : {ws.title} ===")
            rows = []
            for i, row in enumerate(ws.iter_rows(values_only=True)):
                if i >= 200:
                    rows.append("...(lignes suivantes omises)")
                    break
                cells = ["" if c is None else str(c) for c in row]
                rows.append(" | ".join(cells))
            parts.append("\n".join(rows))
        return "\n\n".join(parts)[:30000]
    except Exception as e:
        logger.warning(f"XLSX analyze failed: {e}")
        return ""


def analyze_pptx(data: bytes) -> str:
    """Extrait le texte d'un PPTX (max 60 slides)."""
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(data))
        parts = []
        for i, slide in enumerate(prs.slides):
            if i >= 60:
                break
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for p in shape.text_frame.paragraphs:
                        t = "".join(run.text for run in p.runs).strip()
                        if t:
                            texts.append(t)
            if texts:
                parts.append(f"=== Slide {i + 1} ===\n" + "\n".join(texts))
        return "\n\n".join(parts)[:30000]
    except Exception as e:
        logger.warning(f"PPTX analyze failed: {e}")
        return ""


def analyze_sqlite(data: bytes) -> str:
    """Liste tables + schéma + 10 premières lignes de chaque table d'un SQLite."""
    import sqlite3
    try:
        with tempfile.NamedTemporaryFile(suffix=".sqlite", delete=False) as tf:
            tf.write(data)
            path = tf.name
        try:
            con = sqlite3.connect(path)
            cur = con.cursor()
            parts = []
            tables = [r[0] for r in cur.execute("SELECT name FROM sqlite_master WHERE type='table' ORDER BY name LIMIT 40").fetchall()]
            parts.append(f"Tables ({len(tables)}) : {', '.join(tables) if tables else '(aucune)'}")
            for t in tables[:10]:
                try:
                    schema = cur.execute("SELECT sql FROM sqlite_master WHERE name=?", (t,)).fetchone()
                    parts.append(f"\n--- {t} ---\n{(schema[0] if schema else '')}")
                    rows = cur.execute(f'SELECT * FROM "{t}" LIMIT 10').fetchall()
                    parts.append("Exemples (10 lignes max) :\n" + "\n".join(" | ".join(str(c) for c in r) for r in rows))
                except Exception as te:
                    parts.append(f"(erreur {t}: {te})")
            con.close()
            return "\n".join(parts)[:30000]
        finally:
            try:
                os.unlink(path)
            except Exception:
                pass
    except Exception as e:
        logger.warning(f"SQLITE analyze failed: {e}")
        return ""


# ----------------------------------------------------------------- builders (bytes)
def build_docx_bytes(title: str, sections: List[Dict[str, Any]]) -> bytes:
    from docx import Document
    doc = Document()
    doc.add_heading(title or "Document", 0)
    for sec in sections or []:
        h = (sec.get("heading") or "").strip()
        c = (sec.get("content") or "").strip()
        if h:
            doc.add_heading(h, 1)
        if c:
            for para in c.split("\n\n"):
                doc.add_paragraph(para)
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def build_pdf_bytes(title: str, sections: List[Dict[str, Any]]) -> bytes:
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
    from reportlab.lib.units import cm
    buf = io.BytesIO()
    doc = SimpleDocTemplate(buf, pagesize=A4, title=title or "Document")
    styles = getSampleStyleSheet()
    flow = [Paragraph(title or "Document", styles["Title"]), Spacer(1, 0.6 * cm)]
    for sec in sections or []:
        h = (sec.get("heading") or "").strip()
        c = (sec.get("content") or "").strip()
        if h:
            flow.append(Paragraph(h, styles["Heading1"]))
            flow.append(Spacer(1, 0.2 * cm))
        if c:
            for para in c.split("\n\n"):
                flow.append(Paragraph(para.replace("\n", "<br/>"), styles["BodyText"]))
                flow.append(Spacer(1, 0.2 * cm))
    doc.build(flow)
    return buf.getvalue()


def build_xlsx_bytes(sheets: List[Dict[str, Any]]) -> bytes:
    """sheets = [{"name","headers","rows","formulas"}]"""
    import xlsxwriter
    buf = io.BytesIO()
    wb = xlsxwriter.Workbook(buf, {"in_memory": True})
    header_fmt = wb.add_format({"bold": True, "bg_color": "#E4FF00", "border": 1})
    for idx, sh in enumerate(sheets or [{"name": "Feuille1", "rows": []}]):
        ws = wb.add_worksheet(sh.get("name") or f"Feuille{idx + 1}")
        row_offset = 0
        headers = sh.get("headers") or []
        if headers:
            for c, h in enumerate(headers):
                ws.write(0, c, h, header_fmt)
            row_offset = 1
        for r_i, row in enumerate(sh.get("rows") or []):
            for c_i, val in enumerate(row):
                ws.write(row_offset + r_i, c_i, val)
        for cell, formula in (sh.get("formulas") or {}).items():
            try:
                ws.write_formula(cell, formula)
            except Exception:
                pass
        if headers:
            ws.autofilter(0, 0, row_offset + len(sh.get("rows") or []) - 1, len(headers) - 1)
    wb.close()
    return buf.getvalue()


def build_pptx_bytes(title: str, slides: List[Dict[str, Any]]) -> bytes:
    from pptx import Presentation
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title or "Présentation"
    for s in slides or []:
        layout = prs.slide_layouts[1]
        sl = prs.slides.add_slide(layout)
        sl.shapes.title.text = s.get("title") or ""
        body = sl.placeholders[1]
        tf = body.text_frame
        content = s.get("content") or ""
        lines = content.split("\n")
        tf.text = lines[0] if lines else ""
        for extra in lines[1:]:
            p = tf.add_paragraph()
            p.text = extra
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ----------------------------------------------------------------- Python sandbox
# Persistent session dirs : { session_id: {"dir": path, "last_used": ts } }
_SESSION_DIRS: Dict[str, Dict[str, Any]] = {}
_SESSION_TTL_SEC = 3600  # 1 hour idle → cleanup
_SESSION_MAX = 50  # hard cap on live sessions


def _sandbox_session_dir(session_id: Optional[str]) -> str:
    """Retourne un tmp_dir, PERSISTANT si session_id fourni, sinon éphémère."""
    if not session_id:
        return tempfile.mkdtemp(prefix="cfsandbox_")
    now = time.time()
    # Cleanup expired sessions
    for sid in list(_SESSION_DIRS.keys()):
        if now - _SESSION_DIRS[sid].get("last_used", now) > _SESSION_TTL_SEC:
            try:
                shutil.rmtree(_SESSION_DIRS[sid]["dir"], ignore_errors=True)
            except Exception:
                pass
            _SESSION_DIRS.pop(sid, None)
    # Enforce cap (LRU eviction)
    if session_id not in _SESSION_DIRS and len(_SESSION_DIRS) >= _SESSION_MAX:
        oldest = min(_SESSION_DIRS.items(), key=lambda kv: kv[1].get("last_used", 0))[0]
        try:
            shutil.rmtree(_SESSION_DIRS[oldest]["dir"], ignore_errors=True)
        except Exception:
            pass
        _SESSION_DIRS.pop(oldest, None)
    entry = _SESSION_DIRS.get(session_id)
    if entry and os.path.isdir(entry["dir"]):
        entry["last_used"] = now
        return entry["dir"]
    d = tempfile.mkdtemp(prefix=f"cfrepl_{session_id[:10]}_")
    _SESSION_DIRS[session_id] = {"dir": d, "last_used": now}
    return d


def reset_sandbox_session(session_id: str) -> bool:
    """Supprime le namespace persistant d'une session REPL."""
    entry = _SESSION_DIRS.pop(session_id, None)
    if entry:
        try:
            shutil.rmtree(entry["dir"], ignore_errors=True)
            return True
        except Exception:
            pass
    return False


_SANDBOX_PREAMBLE = (
    "import sys, os\n"
    "os.environ.pop('EMERGENT_LLM_KEY', None)\n"
    "os.environ.pop('RESEND_API_KEY', None)\n"
    "os.environ.pop('MONGO_URL', None)\n"
    "os.environ.pop('DB_NAME', None)\n"
    "os.environ.pop('OLLAMA_BASE_URL', None)\n"
    "sys.setrecursionlimit(1000)\n"
    "# Auto-dump matplotlib figures to ./_figs_/ on plt.show() OR at the end.\n"
    "try:\n"
    "    import matplotlib\n"
    "    matplotlib.use('Agg')\n"
    "    import matplotlib.pyplot as _plt\n"
    "    import atexit as _atexit\n"
    "    os.makedirs('_figs_', exist_ok=True)\n"
    "    _fig_counter = [0]\n"
    "    def _dump_all_figs():\n"
    "        for fnum in _plt.get_fignums():\n"
    "            fig = _plt.figure(fnum)\n"
    "            try:\n"
    "                fig.savefig(f'_figs_/fig_{_fig_counter[0]:03d}.png', bbox_inches='tight', dpi=110)\n"
    "                _fig_counter[0] += 1\n"
    "            except Exception:\n"
    "                pass\n"
    "        _plt.close('all')\n"
    "    _orig_show = _plt.show\n"
    "    def _patched_show(*a, **kw):\n"
    "        _dump_all_figs()\n"
    "    _plt.show = _patched_show\n"
    "    _atexit.register(_dump_all_figs)\n"
    "except Exception:\n"
    "    pass\n"
)


def _empty_result(err: str = "") -> Dict[str, Any]:
    return {"stdout": "", "stderr": err, "exit_code": -1, "timed_out": False, "duration_ms": 0, "images": [], "variables": [], "session_id": None}


async def run_python_sandbox(
    code: str,
    timeout_sec: int = 10,
    session_id: Optional[str] = None,
    files: Optional[List[Dict[str, str]]] = None,
) -> Dict[str, Any]:
    """Exécute du code Python dans un sous-processus avec timeout dur.

    Args:
        code: le code Python à exécuter.
        timeout_sec: timeout hard (1-30 s).
        session_id: si fourni, réutilise le même workspace (variables persistantes
                    entre appels, style REPL/Jupyter). None → exécution éphémère.
        files: liste de fichiers à déposer dans le cwd du script, format
               [{"filename": "data.csv", "data_base64": "..."}]. Max 6 fichiers de 10 Mo.

    Retourne { stdout, stderr, exit_code, timed_out, duration_ms, images[], variables[] }.
    Les figures matplotlib sont automatiquement capturées en base64.
    """
    code = (code or "").strip()
    if not code:
        return _empty_result("Aucun code fourni.")

    persistent = bool(session_id)
    tmp_dir = _sandbox_session_dir(session_id)
    script_path = os.path.join(tmp_dir, "user_script.py")

    # Drop uploaded files (CSV / JSON / etc.) dans le cwd du sandbox.
    if files:
        for f in files[:6]:
            try:
                fname = os.path.basename(f.get("filename") or "")
                if not fname or fname.startswith("."):
                    continue
                raw = base64.b64decode(f.get("data_base64") or "")
                if len(raw) > 10 * 1024 * 1024:
                    continue
                with open(os.path.join(tmp_dir, fname), "wb") as fout:
                    fout.write(raw)
            except Exception as _fx:
                logger.warning(f"sandbox file drop failed: {_fx}")

    # REPL-style namespace persistence via dill.
    state_py = os.path.join(tmp_dir, "_state.pkl")
    has_state = persistent and os.path.isfile(state_py)

    state_restore = (
        "# --- restore state ---\n"
        "try:\n"
        "    import dill as _dill\n"
        "    with open('_state.pkl','rb') as _sf:\n"
        "        _saved = _dill.load(_sf)\n"
        "    for _k,_v in _saved.items():\n"
        "        globals()[_k] = _v\n"
        "except Exception as _re:\n"
        "    pass\n"
    ) if has_state else ""

    state_save = (
        "\n# --- save state (REPL mode) ---\n"
        "try:\n"
        "    import dill as _dill\n"
        "    _to_save = {}\n"
        "    _BAD_TYPES = (type(sys), type(_dill), type(lambda: None).__class__) if False else ()\n"
        "    for _k,_v in list(globals().items()):\n"
        "        if _k.startswith('_') or _k in ('sys','os','dill','matplotlib','plt','atexit'):\n"
        "            continue\n"
        "        try:\n"
        "            _dill.dumps(_v)\n"
        "            _to_save[_k] = _v\n"
        "        except Exception:\n"
        "            pass\n"
        "    with open('_state.pkl','wb') as _sf:\n"
        "        _dill.dump(_to_save, _sf)\n"
        "except Exception:\n"
        "    pass\n"
    ) if persistent else ""

    try:
        with open(script_path, "w", encoding="utf-8") as f:
            f.write(_SANDBOX_PREAMBLE + state_restore + "\n# --- user code ---\n" + code + state_save)

        start = time.monotonic()
        try:
            proc = await asyncio.create_subprocess_exec(
                sys.executable, script_path,
                stdout=asyncio.subprocess.PIPE,
                stderr=asyncio.subprocess.PIPE,
                cwd=tmp_dir,
                env={
                    "PATH": os.environ.get("PATH", "/usr/bin:/usr/local/bin"),
                    "HOME": tmp_dir,
                    "LC_ALL": "C.UTF-8",
                    "LANG": "C.UTF-8",
                    "PYTHONPATH": os.environ.get("PYTHONPATH", ""),
                    "MPLBACKEND": "Agg",
                },
            )
            try:
                stdout_b, stderr_b = await asyncio.wait_for(proc.communicate(), timeout=timeout_sec)
                timed_out = False
            except asyncio.TimeoutError:
                try:
                    proc.kill()
                except Exception:
                    pass
                try:
                    stdout_b, stderr_b = await proc.communicate()
                except Exception:
                    stdout_b, stderr_b = b"", b""
                timed_out = True
            duration_ms = int((time.monotonic() - start) * 1000)
            stdout = (stdout_b or b"").decode("utf-8", errors="replace")[:8000]
            stderr = (stderr_b or b"").decode("utf-8", errors="replace")[:4000]

            # Collect generated images
            images: List[Dict[str, str]] = []
            try:
                figs_dir = os.path.join(tmp_dir, "_figs_")
                if os.path.isdir(figs_dir):
                    for fn in sorted(os.listdir(figs_dir))[:6]:
                        fp = os.path.join(figs_dir, fn)
                        if os.path.isfile(fp) and os.path.getsize(fp) < 4_000_000:
                            with open(fp, "rb") as imf:
                                images.append({
                                    "filename": fn,
                                    "mime_type": "image/png",
                                    "data_base64": base64.b64encode(imf.read()).decode("ascii"),
                                })
                    # Clear _figs_ so next run doesn't re-send them.
                    try:
                        shutil.rmtree(figs_dir, ignore_errors=True)
                    except Exception:
                        pass
                for fn in os.listdir(tmp_dir):
                    if len(images) >= 6:
                        break
                    if fn.lower().endswith((".png", ".jpg", ".jpeg")) and fn not in ("user_script.py",):
                        fp = os.path.join(tmp_dir, fn)
                        if os.path.isfile(fp) and os.path.getsize(fp) < 4_000_000:
                            mt = "image/jpeg" if fn.lower().endswith((".jpg", ".jpeg")) else "image/png"
                            with open(fp, "rb") as imf:
                                images.append({
                                    "filename": fn,
                                    "mime_type": mt,
                                    "data_base64": base64.b64encode(imf.read()).decode("ascii"),
                                })
            except Exception as _imx:
                logger.warning(f"sandbox image collect failed: {_imx}")

            # Report current REPL variables (names + type + short repr).
            variables: List[Dict[str, str]] = []
            if persistent and os.path.isfile(state_py):
                try:
                    import dill as _dill  # type: ignore
                    with open(state_py, "rb") as sf:
                        saved = _dill.load(sf)
                    for k, v in list(saved.items())[:30]:
                        try:
                            r = repr(v)
                            if len(r) > 60:
                                r = r[:57] + "..."
                            variables.append({"name": k, "type": type(v).__name__, "repr": r})
                        except Exception:
                            pass
                except Exception:
                    pass

            return {
                "stdout": stdout,
                "stderr": stderr,
                "exit_code": proc.returncode if proc.returncode is not None else -1,
                "timed_out": timed_out,
                "duration_ms": duration_ms,
                "images": images,
                "variables": variables,
                "session_id": session_id,
            }
        except Exception as e:
            return _empty_result(f"Sandbox error: {e}")
    finally:
        # Only cleanup tmp_dir in ephemeral mode. Persistent dirs are kept.
        if not persistent:
            try:
                shutil.rmtree(tmp_dir, ignore_errors=True)
            except Exception:
                pass
